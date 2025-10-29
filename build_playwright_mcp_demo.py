#!/usr/bin/env python3
"""
Builds a Playwright MCP Demo PowerPoint using your PNG as a full-slide background.

Usage (local):
  pip install -r requirements.txt
  python build_playwright_mcp_demo.py \
    --template-image template.png \
    --out Playwright-MCP-Demo.pptx \
    --presenter "Sheshashayan Papanna" \
    --date "07/11/2025"

This script generates Playwright-MCP-Demo.pptx with consistent slide layout and speaker notes.
"""
import argparse
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_bg_image(slide, image_path, prs):
    slide.shapes.add_picture(
        image_path, 0, 0, width=prs.slide_width, height=prs.slide_height
    )

def add_title(slide, title_text, left=Inches(0.7), top=Inches(0.6), width=Inches(10), height=Inches(1.2)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

def add_subtitle(slide, subtitle_text, left=Inches(0.7), top=Inches(1.6), width=Inches(10), height=Inches(1.5)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)

def add_bullets(slide, bullets, left=Inches(0.7), top=Inches(2.2), width=Inches(10), height=Inches(5.1)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    first = True

    def add_line(text, level=0):
        nonlocal first
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.font.size = Pt(22)

    for item in bullets:
        if isinstance(item, str):
            add_line(item, 0)
        elif isinstance(item, (list, tuple)) and item:
            add_line(item[0], 0)
            for sub in item[1:]:
                add_line(sub, 1)

def add_notes(slide, notes_text):
    if not notes_text:
        return
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.clear()
    tf.paragraphs[0].text = notes_text

def add_content_slide(prs, template_image, title, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg_image(slide, template_image, prs)
    add_title(slide, title)
    add_bullets(slide, bullets)
    add_notes(slide, notes)
    return slide

def add_title_slide(prs, template_image, title, subtitle, presenter, org, when):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg_image(slide, template_image, prs)
    add_title(slide, title)
    subtitle_lines = [subtitle]
    if presenter:
        subtitle_lines.append(presenter)
    if org:
        subtitle_lines.append(org)
    subtitle_lines.append(when)
    add_subtitle(slide, "\n".join(subtitle_lines))
    return slide

def build_prs(args):
    prs = Presentation()
    # Title slide
    add_title_slide(
        prs,
        args.template_image,
        args.title,
        "Definition, Architecture, Components, Uses, Installation, Execution, References",
        args.presenter,
        args.org,
        args.date or str(date.today()),
    )

    slides = [
        ("Agenda", [
            "What is Playwright?",
            "What is MCP (Model Context Protocol)?",
            "Why Playwright + MCP",
            "Architecture",
            "Components",
            "Uses",
            "Install steps",
            "Execute steps (demo flow)",
            "References",
        ], "Set expectations and timing."),
        ("What is Playwright?", [
            "Open-source end-to-end testing framework by Microsoft",
            "Automates Chromium, Firefox, WebKit across Windows, macOS, Linux, CI",
            "Fast, reliable, headless or headed modes",
            "Rich tooling: Codegen, Trace Viewer, HTML reports",
            "Great for E2E, smoke, visual, and API-assisted flows",
        ], "Emphasize cross-browser parity and dev tooling."),
        ("What is MCP (Model Context Protocol)?", [
            "Protocol to safely connect LLMs to external tools and data",
            "Standardizes capabilities (tools/procedures, resources, prompts)",
            "Decouples LLMs from vendor-specific plugins",
            "Enables auditable, permissioned tool usage",
        ], "MCP as a clean contract between LLMs and tools."),
        ("Why Combine Playwright + MCP?", [
            "Natural-language control of browsers via standardized tools",
            "Automate tasks and verifications with LLM guidance",
            "Reproduce flaky flows and triage with traces/logs",
            "Fits CI/CD and guardrails; observability via Playwright artifacts",
        ], "Bridge developer and AI operator workflows."),
        ("Architecture", [
            "LLM Client ↔ MCP Client/Transport ↔ Playwright MCP Server ↔ Playwright Core ↔ Target Web App",
            "Optional: Secrets store, config, logging/telemetry, sandboxing",
        ], "Highlight trust boundaries and policy points."),
        ("Core Components", [
            ["MCP Client:", "Connects LLM to MCP servers", "Presents tools/resources to the model"],
            ["Playwright MCP Server:", "navigate, click, fill, waitFor, screenshot, extract, runTest",
             "Config: browser type, headless, storage state, timeouts"],
            ["Playwright Runtime:", "Browsers, isolation contexts, tracing, reports"],
        ], "Clarify responsibilities and config."),
        ("Common Uses", [
            "E2E and smoke runs via natural-language prompts",
            "Repro/triage regressions with scripts and traces",
            "Data extraction and validation from internal apps",
            "Visual verifications and review screenshots",
            "CI bots that verify 'happy path' before merge",
        ], "Map to concrete workflows."),
        ("Prerequisites", [
            "Node.js 18+ and Git",
            "macOS/Windows/Linux with browser deps",
            "Optional: VS Code and npm/yarn/pnpm",
            "MCP-capable client (e.g., Claude Desktop) for LLM control",
        ], "Note proxies/sandbox considerations."),
        ("Install: Playwright", [
            "npm init -y",
            "npm i -D @playwright/test",
            "npx playwright install --with-deps",
            "Optional: npx playwright test --ui; npx playwright codegen https://example.com",
        ], "Verify browsers downloaded; run a sample test."),
        ("Install: Playwright MCP Server (example)", [
            "git clone <playwright-mcp-server-repo> && cd <repo>",
            "npm i && npm run build",
            "Configure .env: BROWSER, HEADLESS, TIMEOUTS, STORAGE_STATE",
            "npm start (verify local run)",
        ], "Replace with your team’s server repo."),
        ("Configure MCP Client", [
            "Add server in client (e.g., Claude Desktop Tools/Servers)",
            "Command: node dist/index.js; Working dir: server root",
            "Env: HEADLESS=true, BROWSER=chromium, etc.",
            "Validate tools: playwright.navigate, playwright.click, playwright.screenshot",
        ], "Dry run to list available procedures."),
        ("Execute: LLM-Driven Demo", [
            "Start server: npm start",
            "Prompt: Open URL → Click Sign In → Fill creds → Submit",
            "Screenshot and confirm dashboard visible",
            "Collect artifacts: screenshots, traces",
        ], "Keep creds in secrets; show trace viewer."),
        ("Execute: Playwright Test CLI (non-LLM)", [
            "npx playwright test",
            "npx playwright show-report",
            "npx playwright show-trace trace.zip",
        ], "Show parity with scripted flows."),
        ("CI/CD Integration (Example)", [
            "Use official Playwright GitHub Action",
            "Cache browsers for speed",
            "Upload HTML report and trace artifacts",
        ], "Mention parallelization and flaky retry policy."),
        ("Security & Governance", [
            "Least privilege; domain allowlists",
            "Redact secrets/PII; vault-backed env vars",
            "Logging/telemetry with sensitive-data filtering",
            "Sandbox/brokered access for production targets",
        ], "Cover infosec expectations."),
        ("Troubleshooting", [
            "Install browsers: npx playwright install",
            "CI sandbox/headless issues: disable sandbox or run headed locally",
            "Selectors flaky: prefer data-testid and proper waits",
            "MCP connection errors: check path, permissions, env vars",
            "Trace empty: enable tracing around actions",
        ], "Mini-FAQ for demos."),
        ("References", [
            "Playwright Docs: https://playwright.dev",
            "Test Intro: https://playwright.dev/docs/test-intro",
            "CLI: https://playwright.dev/docs/test-cli",
            "CI Guides: https://playwright.dev/docs/ci",
            "Trace Viewer: https://playwright.dev/docs/trace-viewer",
            "MCP Spec: https://modelcontextprotocol.io",
            "Claude + MCP Getting Started: https://docs.anthropic.com/claude/docs/mcp-get-started",
        ], "Add your internal repo link."),
    ]

    for title, bullets, notes in slides:
        add_content_slide(prs, args.template_image, title, bullets, notes)
    return prs

def main():
    ap = argparse.ArgumentParser(description="Generate Playwright MCP Demo PPTX with PNG background")
    ap.add_argument("--template-image", required=True, help="Path to PNG template")
    ap.add_argument("--out", default="Playwright-MCP-Demo.pptx")
    ap.add_argument("--title", default="Playwright MCP Demo")
    ap.add_argument("--presenter", default="Sheshashayan Papanna")
    ap.add_argument("--org", default="")
    ap.add_argument("--date", default="07/11/2025")
    args = ap.parse_args()

    prs = build_prs(args)
    prs.save(args.out)
    print(f"Wrote {args.out}")


if __name__ == "__main__":
    main()